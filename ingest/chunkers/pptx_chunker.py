import hashlib
import os
from typing import Iterator

import tiktoken
from pptx import Presentation

from .base import Chunk

CHUNK_SIZE = 800
CHUNK_OVERLAP = 150
ENCODER = tiktoken.get_encoding("cl100k_base")


def _token_count(text: str) -> int:
    return len(ENCODER.encode(text))


def _split_by_tokens(text: str, chunk_size: int, overlap: int) -> list[str]:
    tokens = ENCODER.encode(text)
    chunks = []
    start = 0
    while start < len(tokens):
        end = min(start + chunk_size, len(tokens))
        chunks.append(ENCODER.decode(tokens[start:end]))
        if end == len(tokens):
            break
        start += chunk_size - overlap
    return chunks


def _extract_slide_text(slide) -> tuple[str, str]:
    """Retourne (titre, corps) pour une slide."""
    title = ""
    body_parts = []

    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        text = shape.text_frame.text.strip()
        if not text:
            continue
        if shape.is_placeholder and shape.placeholder_format.idx == 0:
            title = text
        else:
            body_parts.append(text)

    return title, "\n\n".join(body_parts)


class PPTXChunker:
    def chunk_file(self, file_path: str) -> Iterator[Chunk]:
        with open(file_path, "rb") as f:
            raw_bytes = f.read()

        file_hash = hashlib.sha256(raw_bytes).hexdigest()
        source_file = os.path.basename(file_path)
        prs = Presentation(file_path)

        chunk_index = 0
        for slide_num, slide in enumerate(prs.slides, start=1):
            title, body = _extract_slide_text(slide)

            if not title and not body:
                continue

            content = f"{title}\n\n{body}".strip() if title and body else (title or body)

            if _token_count(content) <= CHUNK_SIZE:
                yield Chunk(
                    content=content,
                    source_file=source_file,
                    page_number=slide_num,
                    chunk_index=chunk_index,
                    doc_type="pptx",
                    section=title,
                    title=title or source_file,
                    file_hash=file_hash,
                )
                chunk_index += 1
            else:
                for sub in _split_by_tokens(content, CHUNK_SIZE, CHUNK_OVERLAP):
                    yield Chunk(
                        content=sub,
                        source_file=source_file,
                        page_number=slide_num,
                        chunk_index=chunk_index,
                        doc_type="pptx",
                        section=title,
                        title=title or source_file,
                        file_hash=file_hash,
                    )
                    chunk_index += 1
